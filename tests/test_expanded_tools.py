from io import BytesIO

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas


def _docx_bytes() -> BytesIO:
    buffer = BytesIO()
    document = Document()
    document.add_heading("Knowledge Title", level=1)
    document.add_paragraph("Agent content.")
    document.save(buffer)
    buffer.seek(0)
    return buffer


def _xlsx_bytes() -> BytesIO:
    buffer = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["name", "value"])
    sheet.append(["alpha", 1])
    workbook.save(buffer)
    buffer.seek(0)
    return buffer


def _pptx_bytes() -> BytesIO:
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Progress summary"
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def _pdf_bytes() -> BytesIO:
    buffer = BytesIO()
    document = canvas.Canvas(buffer)
    document.drawString(72, 720, "PDF knowledge text")
    document.save()
    buffer.seek(0)
    return buffer


def test_tool_list_contains_expanded_agent_capabilities(client, headers):
    response = client.get("/v1/tools", headers=headers)

    assert response.status_code == 200
    names = {tool["name"] for tool in response.json()["tools"]}
    assert {
        "document_parse",
        "pdf_parse",
        "word_parse",
        "excel_parse",
        "ppt_parse",
        "html_parse",
        "txt_parse",
        "md_parse",
        "img_parse",
        "word_to_md",
        "pdf_to_md",
        "image_ocr",
        "pdf_ocr",
        "chunks_data_parse",
        "bm25_preprocess",
        "http_request",
        "send_email_smtp",
        "jira_search_issues",
    }.issubset(names)


def test_process_run_parses_uploaded_markdown(client, headers):
    response = client.post(
        "/v1/process/run",
        headers=headers,
        data={"tool_name": "document_parse"},
        files={"file": ("knowledge.md", BytesIO(b"# Title\n\nBody"), "text/markdown")},
    )

    assert response.status_code == 200
    data = response.json()["result"]["data"]
    assert data["filename"] == "knowledge.md"
    assert data["section_count"] == 1
    assert "# Title" in data["sections"][0]["content"]


def test_process_run_parses_excel_and_ppt_alias_tools(client, headers):
    excel = client.post(
        "/v1/process/run",
        headers=headers,
        data={"tool_name": "excel_parse"},
        files={"file": ("knowledge.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
    )
    assert excel.status_code == 200
    assert excel.json()["result"]["data"]["sections"][0]["sheet"] == "Data"
    assert "alpha" in excel.json()["result"]["data"]["sections"][0]["content"]

    ppt = client.post(
        "/v1/process/run",
        headers=headers,
        data={"tool_name": "ppt_parse"},
        files={"file": ("slides.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert ppt.status_code == 200
    assert "Quarterly Review" in ppt.json()["result"]["data"]["sections"][0]["content"]


def test_word_to_markdown_supports_direct_download(client, headers):
    response = client.post(
        "/v1/convert/download",
        headers=headers,
        data={"tool_name": "word_to_md", "output_filename": "knowledge.md"},
        files={
            "file": (
                "knowledge.docx",
                _docx_bytes(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )

    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/markdown")
    assert b"Knowledge Title" in response.content


def test_pdf_to_markdown_supports_direct_download(client, headers):
    response = client.post(
        "/v1/convert/download",
        headers=headers,
        data={"tool_name": "pdf_to_md", "output_filename": "knowledge.md"},
        files={"file": ("knowledge.pdf", _pdf_bytes(), "application/pdf")},
    )

    assert response.status_code == 200
    assert b"PDF knowledge text" in response.content


def test_text_chunk_and_bm25_preprocess(client, headers):
    chunk_response = client.post(
        "/v1/tools/chunks_data_parse/run",
        headers=headers,
        json={"input": {"text": "abcdefghij", "chunk_size": 5, "overlap": 1}},
    )
    assert chunk_response.status_code == 200
    assert chunk_response.json()["result"]["data"]["count"] == 3

    bm25_response = client.post(
        "/v1/tools/bm25_preprocess/run",
        headers=headers,
        json={"input": {"text": "AI Agent 文档"}},
    )
    assert bm25_response.status_code == 200
    assert "ai" in bm25_response.json()["result"]["data"]["tokens"]


def test_process_route_is_described_in_openapi(client):
    spec = client.get("/openapi.json").json()

    assert "/v1/process/run" in spec["paths"]
    assert "OCR" in spec["paths"]["/v1/process/run"]["post"]["summary"]
