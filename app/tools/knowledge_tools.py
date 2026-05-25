import base64
import io
import json
import re
from pathlib import Path
from typing import Any

import markdownify
from bs4 import BeautifulSoup
from docx import Document
from PIL import Image, ImageDraw
from pypdf import PdfReader

from app.storage.base import Storage
from app.tools.base import BaseTool, ToolExecutionError, ToolResult
from app.tools.utils import detect_content_type, require_field


def _stored_file(input_data: dict[str, Any], storage: Storage) -> tuple[Path, dict[str, Any]]:
    file_id = require_field(input_data, "file_id")
    meta = storage.get(file_id)
    return Path(meta["path"]), meta


def _docx_text(path: Path) -> str:
    document = Document(str(path))
    lines: list[str] = []
    for paragraph in document.paragraphs:
        style = paragraph.style.name.lower() if paragraph.style else ""
        text = paragraph.text.strip()
        if not text:
            continue
        if style.startswith("heading"):
            level = re.search(r"\d+", style)
            lines.append(f"{'#' * int(level.group() if level else 1)} {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines)


def _pdf_pages(path: Path) -> list[dict[str, Any]]:
    reader = PdfReader(str(path))
    return [{"page": index + 1, "content": page.extract_text() or ""} for index, page in enumerate(reader.pages)]


class WordToMarkdownTool(BaseTool):
    name = "word_to_md"
    description = "Convert DOCX content to Markdown text or a downloadable Markdown file."
    input_schema = {
        "type": "object",
        "required": ["file_id"],
        "properties": {
            "file_id": {"type": "string", "description": "已上传的 DOCX 文件 ID。", "examples": ["report.docx"]},
            "output_filename": {"type": "string", "default": "document.md", "examples": ["report.md"]},
        },
    }

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        source, meta = _stored_file(input_data, storage)
        if source.suffix.lower() != ".docx":
            raise ToolExecutionError("invalid_input", "word_to_md currently supports .docx files.")
        markdown = _docx_text(source)
        filename = input_data.get("output_filename") or f"{Path(meta['filename']).stem}.md"
        stored = storage.save_bytes(markdown.encode("utf-8"), filename, "text/markdown")
        return ToolResult(
            type="file",
            data={"markdown": markdown},
            result_file_id=stored["file_id"],
            filename=stored["filename"],
            content_type=stored["content_type"],
        )


class PdfToMarkdownTool(BaseTool):
    name = "pdf_to_md"
    description = "Extract PDF text page by page and save it as a Markdown file."
    input_schema = {
        "type": "object",
        "required": ["file_id"],
        "properties": {
            "file_id": {"type": "string", "description": "已上传的 PDF 文件 ID。", "examples": ["report.pdf"]},
            "output_filename": {"type": "string", "default": "document.md", "examples": ["report.md"]},
        },
    }

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        source, meta = _stored_file(input_data, storage)
        if source.suffix.lower() != ".pdf":
            raise ToolExecutionError("invalid_input", "pdf_to_md expects a .pdf file.")
        pages = _pdf_pages(source)
        markdown = "\n\n".join(f"## Page {item['page']}\n\n{item['content'].strip()}" for item in pages)
        filename = input_data.get("output_filename") or f"{Path(meta['filename']).stem}.md"
        stored = storage.save_bytes(markdown.encode("utf-8"), filename, "text/markdown")
        return ToolResult(
            type="file",
            data={"page_count": len(pages)},
            result_file_id=stored["file_id"],
            filename=stored["filename"],
            content_type=stored["content_type"],
        )


class DocumentParseTool(BaseTool):
    name = "document_parse"
    description = "Parse common document files into page/sheet/slide-oriented content for RAG ingestion."
    input_schema = {
        "type": "object",
        "required": ["file_id"],
        "properties": {
            "file_id": {
                "type": "string",
                "description": "支持 PDF, DOCX, XLSX, PPTX, HTML, TXT, Markdown。",
                "examples": ["knowledge.pdf"],
            },
            "encoding": {"type": "string", "default": "utf-8", "examples": ["utf-8"]},
        },
    }
    timeout_seconds = 600

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        source, meta = _stored_file(input_data, storage)
        suffix = source.suffix.lower()
        content: list[dict[str, Any]]
        if suffix == ".pdf":
            content = _pdf_pages(source)
        elif suffix == ".docx":
            content = [{"page": 1, "content": _docx_text(source)}]
        elif suffix in {".html", ".htm"}:
            html = source.read_text(encoding=input_data.get("encoding", "utf-8"))
            content = [{"page": 1, "content": markdownify.markdownify(html, heading_style="ATX")}]
        elif suffix in {".txt", ".md", ".json", ".yaml", ".yml"}:
            content = [{"page": 1, "content": source.read_text(encoding=input_data.get("encoding", "utf-8"))}]
        elif suffix == ".xlsx":
            content = _parse_xlsx(source)
        elif suffix == ".pptx":
            content = _parse_pptx(source)
        else:
            raise ToolExecutionError("unsupported_file_type", f"document_parse does not support {suffix}.")
        return ToolResult(
            type="json",
            data={
                "filename": meta["filename"],
                "content_type": detect_content_type(meta["filename"]),
                "sections": content,
                "section_count": len(content),
            },
        )


class _TypedDocumentParseTool(DocumentParseTool):
    supported_suffixes: set[str]

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        source, _ = _stored_file(input_data, storage)
        if source.suffix.lower() not in self.supported_suffixes:
            allowed = ", ".join(sorted(self.supported_suffixes))
            raise ToolExecutionError("invalid_input", f"{self.name} expects one of: {allowed}.")
        return super().execute(input_data, storage)


class PdfParseTool(_TypedDocumentParseTool):
    name = "pdf_parse"
    description = "Parse a PDF into page-oriented text sections for knowledge ingestion."
    supported_suffixes = {".pdf"}


class WordParseTool(_TypedDocumentParseTool):
    name = "word_parse"
    description = "Parse a DOCX document into structured text sections for knowledge ingestion."
    supported_suffixes = {".docx"}


class ExcelParseTool(_TypedDocumentParseTool):
    name = "excel_parse"
    description = "Parse an XLSX workbook into sheet-oriented content for knowledge ingestion."
    supported_suffixes = {".xlsx"}


class PptParseTool(_TypedDocumentParseTool):
    name = "ppt_parse"
    description = "Parse a PPTX presentation into slide-oriented content for knowledge ingestion."
    supported_suffixes = {".pptx"}


class HtmlParseTool(_TypedDocumentParseTool):
    name = "html_parse"
    description = "Parse an HTML document into Markdown content for knowledge ingestion."
    supported_suffixes = {".html", ".htm"}


class TxtParseTool(_TypedDocumentParseTool):
    name = "txt_parse"
    description = "Parse a plain text document for knowledge ingestion."
    supported_suffixes = {".txt"}


class MarkdownParseTool(_TypedDocumentParseTool):
    name = "md_parse"
    description = "Parse a Markdown document for knowledge ingestion."
    supported_suffixes = {".md"}


def _parse_xlsx(path: Path) -> list[dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ToolExecutionError("missing_dependency", "openpyxl is required for XLSX parsing.") from exc
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections = []
    for sheet in workbook.worksheets:
        rows = [["" if value is None else str(value) for value in row] for row in sheet.iter_rows(values_only=True)]
        sections.append({"sheet": sheet.title, "rows": rows, "content": "\n".join("\t".join(row) for row in rows)})
    return sections


def _parse_pptx(path: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ToolExecutionError("missing_dependency", "python-pptx is required for PPTX parsing.") from exc
    presentation = Presentation(str(path))
    sections = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
        sections.append({"slide": index, "content": text})
    return sections


class TextChunkTool(BaseTool):
    name = "chunks_data_parse"
    description = "Split text or parsed document sections into overlapping chunks for RAG indexing."
    input_schema = {
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "需要切片的原始文本。"},
            "sections": {"type": "array", "description": "document_parse 输出的 sections。"},
            "chunk_size": {"type": "integer", "default": 1000, "minimum": 1},
            "overlap": {"type": "integer", "default": 100, "minimum": 0},
        },
    }

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        text = input_data.get("text")
        if not text and isinstance(input_data.get("sections"), list):
            text = "\n\n".join(str(item.get("content", "")) for item in input_data["sections"])
        if not text:
            raise ToolExecutionError("invalid_input", "Provide text or sections.")
        chunk_size = int(input_data.get("chunk_size", 1000))
        overlap = int(input_data.get("overlap", 100))
        if chunk_size <= 0 or overlap < 0 or overlap >= chunk_size:
            raise ToolExecutionError("invalid_input", "chunk_size must be positive and overlap must be smaller than chunk_size.")
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunks.append({"index": len(chunks), "start": start, "end": end, "content": text[start:end]})
            if end == len(text):
                break
            start = end - overlap
        return ToolResult(type="json", data={"chunks": chunks, "count": len(chunks)})


class Bm25PreprocessTool(BaseTool):
    name = "bm25_preprocess"
    description = "Normalize text and produce searchable tokens for BM25 indexing."
    input_schema = {
        "type": "object",
        "required": ["text"],
        "properties": {"text": {"type": "string", "description": "待预处理文本。", "examples": ["AI agent 文档解析 API"]}},
    }

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        text = require_field(input_data, "text")
        normalized = re.sub(r"\s+", " ", text).strip().lower()
        tokens = re.findall(r"[\u4e00-\u9fff]|[a-z0-9_]+", normalized)
        return ToolResult(type="json", data={"normalized_text": normalized, "tokens": tokens})


class ImageOcrTool(BaseTool):
    name = "image_ocr"
    description = "Perform OCR on an uploaded image or base64 image using local Tesseract."
    input_schema = {
        "type": "object",
        "properties": {
            "file_id": {"type": "string", "description": "已上传图片文件 ID。优先于 image_base64。"},
            "image_base64": {"type": "string", "description": "没有上传文件时可传图片 base64。"},
            "language": {"type": "string", "default": "eng", "examples": ["chi_sim+eng"]},
        },
    }
    timeout_seconds = 600

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        image = _load_image(input_data, storage)
        pytesseract = _pytesseract()
        try:
            text = pytesseract.image_to_string(image, lang=str(input_data.get("language", "eng")))
        except Exception as exc:
            raise ToolExecutionError("ocr_failed", "Tesseract OCR execution failed.", {"error": str(exc)}) from exc
        return ToolResult(type="text", data={"text": text})


class ImageParseTool(ImageOcrTool):
    name = "img_parse"
    description = "Parse an uploaded image by OCR and return extracted text for knowledge ingestion."


class PdfOcrTool(BaseTool):
    name = "pdf_ocr"
    description = "Render PDF pages as images and OCR each page using local Tesseract."
    input_schema = {
        "type": "object",
        "required": ["file_id"],
        "properties": {
            "file_id": {"type": "string", "description": "已上传 PDF 文件 ID。"},
            "language": {"type": "string", "default": "eng", "examples": ["chi_sim+eng"]},
            "dpi": {"type": "integer", "default": 200},
        },
    }
    timeout_seconds = 1200

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        source, _ = _stored_file(input_data, storage)
        if source.suffix.lower() != ".pdf":
            raise ToolExecutionError("invalid_input", "pdf_ocr expects a .pdf file.")
        try:
            import fitz
        except ImportError as exc:
            raise ToolExecutionError("missing_dependency", "PyMuPDF is required for PDF OCR.") from exc
        pytesseract = _pytesseract()
        dpi = int(input_data.get("dpi", 200))
        pages = []
        try:
            with fitz.open(str(source)) as document:
                matrix = fitz.Matrix(dpi / 72, dpi / 72)
                for index, page in enumerate(document, start=1):
                    pixmap = page.get_pixmap(matrix=matrix)
                    image = Image.open(io.BytesIO(pixmap.tobytes("png")))
                    text = pytesseract.image_to_string(image, lang=str(input_data.get("language", "eng")))
                    pages.append({"page": index, "content": text})
        except Exception as exc:
            raise ToolExecutionError("ocr_failed", "PDF OCR execution failed.", {"error": str(exc)}) from exc
        return ToolResult(type="json", data={"pages": pages, "page_count": len(pages)})


class ImageAnnotateLayoutTool(BaseTool):
    name = "image_annotate_layout"
    description = "OCR an image and return detected text boxes plus an annotated image."
    input_schema = ImageOcrTool.input_schema
    timeout_seconds = 600

    def execute(self, input_data: dict[str, Any], storage: Storage) -> ToolResult:
        image = _load_image(input_data, storage).convert("RGB")
        pytesseract = _pytesseract()
        try:
            data = pytesseract.image_to_data(
                image,
                lang=str(input_data.get("language", "eng")),
                output_type=pytesseract.Output.DICT,
            )
        except Exception as exc:
            raise ToolExecutionError("ocr_failed", "Tesseract layout detection failed.", {"error": str(exc)}) from exc
        draw = ImageDraw.Draw(image)
        boxes = []
        for index, text in enumerate(data["text"]):
            if not text.strip():
                continue
            box = {
                "text": text,
                "left": data["left"][index],
                "top": data["top"][index],
                "width": data["width"][index],
                "height": data["height"][index],
                "confidence": data["conf"][index],
            }
            boxes.append(box)
            draw.rectangle(
                (box["left"], box["top"], box["left"] + box["width"], box["top"] + box["height"]),
                outline="red",
                width=2,
            )
        output = storage.make_temp_dir() / "annotated.png"
        image.save(output)
        stored = storage.save_path(output, "annotated.png", "image/png")
        return ToolResult(
            type="file",
            data={"boxes": boxes},
            result_file_id=stored["file_id"],
            filename=stored["filename"],
            content_type=stored["content_type"],
        )


def _pytesseract():
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        return pytesseract
    except Exception as exc:
        raise ToolExecutionError(
            "missing_dependency",
            "Tesseract OCR is required. Install Tesseract and make it available on PATH.",
        ) from exc


def _load_image(input_data: dict[str, Any], storage: Storage) -> Image.Image:
    if input_data.get("file_id"):
        path = storage.path_for(str(input_data["file_id"]))
        return Image.open(path)
    if input_data.get("image_base64"):
        try:
            encoded = str(input_data["image_base64"]).split(",", 1)[-1]
            return Image.open(io.BytesIO(base64.b64decode(encoded)))
        except Exception as exc:
            raise ToolExecutionError("invalid_input", "image_base64 is not a readable image.") from exc
    raise ToolExecutionError("invalid_input", "Provide file_id or image_base64.")
